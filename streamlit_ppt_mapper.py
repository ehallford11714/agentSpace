import streamlit as st
import pandas as pd
from io import BytesIO
from pptx import Presentation
import json
import openai


def call_llm(api_key: str, instructions: str, dataframes: list) -> str:
    """Use OpenAI to interpret instructions and return actions in JSON."""
    openai.api_key = api_key

    csv_preview = "\n".join(
        f"CSV {i+1}:\n{df.head().to_csv(index=False)}" for i, df in enumerate(dataframes)
    )
    prompt = (
        "You are an assistant that maps CSV data to a PowerPoint template.\n"
        f"Instructions:\n{instructions}\n"
        "Here are previews of the CSV files:\n" + csv_preview + "\n"
        "Respond with a JSON array of actions. Each action should have "
        "'slide_index', 'placeholder', and 'text'."
    )

    response = openai.chat.completions.create(
        model="gpt-4",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
    )
    return response.choices[0].message.content.strip()


def apply_actions(prs: Presentation, actions_json: str) -> Presentation:
    """Apply LLM actions to the presentation."""
    actions = json.loads(actions_json)
    for action in actions:
        slide_idx = action.get("slide_index")
        placeholder = action.get("placeholder")
        text = action.get("text")
        if slide_idx is None or placeholder is None:
            continue
        try:
            slide = prs.slides[slide_idx]
        except IndexError:
            continue
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.name == placeholder:
                shape.text = text
    return prs


def main() -> None:
    st.title("PowerPoint Data Mapper")
    st.write(
        "Upload a PPTX template, a text file with mapping instructions, and one or more CSV files."
    )

    pptx_file = st.file_uploader("PPTX Template", type=["pptx"])
    instructions_file = st.file_uploader("Instructions (txt)", type=["txt"])
    csv_files = st.file_uploader(
        "CSV Files", type=["csv"], accept_multiple_files=True
    )
    api_key = st.text_input("OpenAI API Key", type="password")

    if (
        st.button("Generate Presentation")
        and pptx_file
        and instructions_file
        and csv_files
        and api_key
    ):
        try:
            instructions = instructions_file.read().decode("utf-8")
        except Exception as e:
            st.error(f"Failed to read instructions: {e}")
            return

        dataframes = []
        for f in csv_files:
            try:
                df = pd.read_csv(f)
                dataframes.append(df)
            except Exception as e:
                st.error(f"Failed to read CSV {f.name}: {e}")
                return

        st.info("Contacting LLM...")
        try:
            actions_json = call_llm(api_key, instructions, dataframes)
            st.text_area("LLM Response", actions_json, height=200)
        except Exception as e:
            st.error(f"LLM call failed: {e}")
            return

        prs = Presentation(BytesIO(pptx_file.read()))
        try:
            prs = apply_actions(prs, actions_json)
            output = BytesIO()
            prs.save(output)
            st.success("Presentation generated!")
            st.download_button(
                label="Download Updated PPTX",
                data=output.getvalue(),
                file_name="updated.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        except Exception as e:
            st.error(f"Failed to apply actions: {e}")


if __name__ == "__main__":
    main()
